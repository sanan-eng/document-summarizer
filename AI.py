#
# import os
# import pdfplumber
# from docx import Document
# from pptx import Presentation
# from groq import Groq
# from langchain_text_splitters import RecursiveCharacterTextSplitter
#
#
# # TEXT EXTRACTION FUNCTIONS
#
# def extract_text(file_path: str) -> str:
#     """
#     Extract text from PDF, DOCX, or PPTX files
#
#     Args:
#         file_path (str): Path to the file
#
#     Returns:
#         str: Extracted text
#
#     Raises:
#         FileNotFoundError: If file doesn't exist
#         ValueError: If file format is unsupported
#     """
#     file_path = file_path.strip().strip('"').strip("'")
#
#     if not os.path.exists(file_path):
#         raise FileNotFoundError(f"File not found: {file_path}")
#
#     file_type = file_path.split(".")[-1].lower()
#
#     if file_type == "pdf":
#         return extract_pdf_text(file_path)
#     elif file_type == "docx":
#         return extract_docx_text(file_path)
#     elif file_type == "pptx":
#         return extract_pptx_text(file_path)
#     else:
#         raise ValueError(f"Unsupported file format: .{file_type}")
#
#
# def extract_pdf_text(file_path: str) -> str:
#     """Extract text from PDF files"""
#     extracted_text = ""
#     with pdfplumber.open(file_path) as pdf:
#         for page in pdf.pages:
#             text = page.extract_text()
#             if text:
#                 extracted_text += text + "\n"
#     return extracted_text.strip()
#
#
# def extract_docx_text(file_path: str) -> str:
#     """Extract text from DOCX files"""
#     doc = Document(file_path)
#     return "\n".join([para.text for para in doc.paragraphs]).strip()
#
#
# def extract_pptx_text(file_path: str) -> str:
#     """Extract text from PPTX files"""
#     prs = Presentation(file_path)
#     text = ""
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if hasattr(shape, "text"):
#                 text += shape.text + "\n"
#     return text.strip()
#
#
# # GROQ SETUP - Using environment variable for API key
# # GROQ SETUP - API key directly in code
# GROQ_API_KEY ="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"
#
# client = Groq(api_key=GROQ_API_KEY)
#
#
# splitter = RecursiveCharacterTextSplitter(
#     chunk_size=2000,
#     chunk_overlap=200
# )
#
# # Consistent system prompt as per requirements
# SYSTEM_PROMPT = "You are an expert document summarization AI that produces accurate, coherent, and well-structured summaries without losing important meaning."
#
#
# def _summarize(text: str, prompt: str) -> str:
#     """
#     Internal function to summarize text using Groq LLM
#
#     Args:
#         text (str): Text to summarize
#         prompt (str): User prompt for summarization
#
#     Returns:
#         str: Generated summary
#     """
#     if not text or not text.strip():
#         raise ValueError("Text cannot be empty")
#
#     chunks = splitter.split_text(text)
#     final_summary = []
#
#     for i, chunk in enumerate(chunks):
#         response = client.chat.completions.create(
#             model="llama-3.1-8b-instant",
#             messages=[
#                 {"role": "system", "content": SYSTEM_PROMPT},
#                 {"role": "user", "content": f"{prompt}\n\nText to summarize:\n{chunk}"}
#             ],
#             temperature=0.3,
#             max_tokens=1000
#         )
#         final_summary.append(response.choices[0].message.content.strip())
#
#     return "\n\n".join(final_summary)
#
#
# # 3 SUMMARY FUNCTIONS with dynamic prompts
#
# def short_summary(text: str) -> str:
#     """Generate very concise bullet points summary"""
#     user_prompt = """Create a very concise collective summary of the ENTIRE document using 7-10 bullet points.
#     Focus only on the most important key information from all parts of the document.
#     Ensure the summary flows as a single, coherent overview."""
#     return _summarize(text, user_prompt)
#
#
# def medium_summary(text: str) -> str:
#     """Generate clear paragraph-style summary"""
#     user_prompt = "Write a clear, moderately detailed summary in paragraph form (150-250 words). Maintain logical flow and highlight main points."
#     return _summarize(text, user_prompt)
#
#
# def long_summary(text: str) -> str:
#     """Generate structured summary with headings and explanations"""
#     user_prompt = "Create a detailed, structured summary with clear headings and explanations. Cover all important aspects comprehensively."
#     return _summarize(text, user_prompt)

import os
import pdfplumber
from docx import Document
from pptx import Presentation
from groq import Groq
from langchain_text_splitters import RecursiveCharacterTextSplitter


# ================= TEXT EXTRACTION =================

def extract_text(file_path: str) -> str:
    file_path = file_path.strip().strip('"').strip("'")

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    file_type = file_path.split(".")[-1].lower()

    if file_type == "pdf":
        return extract_pdf_text(file_path)
    elif file_type == "docx":
        return extract_docx_text(file_path)
    elif file_type == "pptx":
        return extract_pptx_text(file_path)
    elif file_type == "txt":
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    else:
        raise ValueError(f"Unsupported file format: .{file_type}")


def extract_pdf_text(file_path: str) -> str:
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    return text.strip()


def extract_docx_text(file_path: str) -> str:
    doc = Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs).strip()


def extract_pptx_text(file_path: str) -> str:
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()


# ================= GROQ SETUP =================
# ⚠️ API KEY DIRECT PASTE (DEV PURPOSE)


GROQ_API_KEY = os.getenv("GROQ_API_KEY")

client = Groq(api_key=GROQ_API_KEY)

SYSTEM_PROMPT = (
    "You are an expert document analysis and summarization AI. "
    "You produce accurate, coherent, and high-quality summaries."
)

splitter = RecursiveCharacterTextSplitter(
    chunk_size=2000,
    chunk_overlap=200
)


# ================= PREPROCESSING (LLM POWER) =================

def preprocess_text_with_llm(raw_text: str) -> str:
    response = client.chat.completions.create(
        model="llama-3.1-8b-instant",
        messages=[
            {
                "role": "system",
                "content": (
                    "You are an expert document preprocessing AI. "
                    "Your task is to clean and normalize extracted document text."
                )
            },
            {
                "role": "user",
                "content": f"""
INSTRUCTIONS:
- Remove headers, footers, page numbers
- Remove duplicated or repeated content
- Fix broken sentences and paragraphs
- Convert bullet points into complete sentences
- Remove formatting noise and symbols
- Preserve all important information
- Maintain logical flow
- DO NOT summarize

TEXT:
{raw_text}
"""
            }
        ],
        temperature=0.2,
        max_tokens=3000
    )

    return response.choices[0].message.content.strip()


# ================= CORE SUMMARIZATION =================

def _summarize(text: str, final_prompt: str) -> str:
    if not text.strip():
        raise ValueError("Text cannot be empty")

    # 1️⃣ PREPROCESS FULL DOCUMENT
    cleaned_text = preprocess_text_with_llm(text)

    # 2️⃣ SPLIT CLEAN TEXT
    chunks = splitter.split_text(cleaned_text)
    chunk_summaries = []

    # 3️⃣ CHUNK-LEVEL SUMMARIES
    for chunk in chunks:
        response = client.chat.completions.create(
            model="llama-3.1-8b-instant",
            messages=[
                {"role": "system", "content": SYSTEM_PROMPT},
                {
                    "role": "user",
                    "content": f"""
Summarize this section clearly and accurately.
Focus only on meaningful ideas and key points.

TEXT:
{chunk}
"""
                }
            ],
            temperature=0.3,
            max_tokens=700
        )

        chunk_summaries.append(
            response.choices[0].message.content.strip()
        )

    # 4️⃣ FINAL GLOBAL SUMMARY
    combined_summary = "\n\n".join(chunk_summaries)

    final_response = client.chat.completions.create(
        model="llama-3.1-8b-instant",
        messages=[
            {"role": "system", "content": SYSTEM_PROMPT},
            {
                "role": "user",
                "content": f"""
{final_prompt}

DOCUMENT CONTENT:
{combined_summary}
"""
            }
        ],
        temperature=0.3,
        max_tokens=1200
    )

    return final_response.choices[0].message.content.strip()


# ================= SUMMARY TYPES =================

def short_summary(text: str) -> str:
    prompt = """
Create an executive-level summary using 7–10 concise bullet points.
Focus on core ideas, insights, and conclusions.
Avoid repetition and fluff.
"""
    return _summarize(text, prompt)


def medium_summary(text: str) -> str:
    prompt = """
Write a professional and coherent summary of 180–250 words.
Maintain logical flow.
Cover all major concepts, arguments, and conclusions.
"""
    return _summarize(text, prompt)


def long_summary(text: str) -> str:
    prompt = """
Create a detailed, structured summary with clear headings:
- Introduction
- Main Sections
- Key Findings / Arguments
- Conclusion

Explain ideas clearly and comprehensively.
Avoid redundancy.
"""
    return _summarize(text, prompt)
